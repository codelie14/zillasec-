import io
import pandas as pd
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches

import models

def create_pdf_report(analysis: models.Analysis) -> io.BytesIO:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []
    
    res = analysis.analysis_result

    story.append(Paragraph(f"Analysis Report for: {res['metadata']['fichier']}", styles['h1']))
    story.append(Paragraph(f"Date: {res['metadata']['date_analyse']}", styles['h2']))
    story.append(Spacer(1, 12))

    # --- Stats ---
    story.append(Paragraph("Statistics", styles['h2']))
    stats_data = [[k.replace('_', ' ').title(), v] for k, v in res['statistiques'].items()]
    stats_table = Table(stats_data, colWidths=[200, 100])
    stats_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(stats_table)
    story.append(Spacer(1, 12))

    # --- DB Verification & Alerts ---
    story.append(Paragraph("Verification & Alerts", styles['h2']))
    db_alerts_data = [
        ["Present in DB", res['verification_bd']['comptes_presents']],
        ["Absent from DB", res['verification_bd']['comptes_absents']],
        ["Disabled Admins", res['alertes']['admin_desactives']],
        ["Disabled Sensitive Access", res['alertes']['acces_sensibles_desactives']],
    ]
    db_alerts_table = Table(db_alerts_data, colWidths=[200, 100])
    db_alerts_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(db_alerts_table)
    story.append(Spacer(1, 12))

    # --- Account Details ---
    story.append(Paragraph("Account Details", styles['h2']))
    account_data = [list(res['details_comptes'][0].keys())] # Headers
    for account in res['details_comptes']:
        account_data.append(list(str(v) for v in account.values()))

    account_table = Table(account_data)
    account_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(account_table)

    doc.build(story)
    buffer.seek(0)
    return buffer

def create_excel_report(analysis: models.Analysis) -> io.BytesIO:
    buffer = io.BytesIO()
    res = analysis.analysis_result
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        pd.DataFrame.from_dict(res['statistiques'], orient='index', columns=['Value']).to_excel(writer, sheet_name='Statistics')
        pd.DataFrame.from_dict(res['verification_bd'], orient='index', columns=['Value']).to_excel(writer, sheet_name='DB Verification')
        pd.DataFrame.from_dict(res['alertes'], orient='index', columns=['Value']).to_excel(writer, sheet_name='Alerts')
        pd.DataFrame(res['details_comptes']).to_excel(writer, sheet_name='Account Details', index=False)
    buffer.seek(0)
    return buffer

def create_pptx_report(analysis: models.Analysis) -> io.BytesIO:
    buffer = io.BytesIO()
    prs = Presentation()
    res = analysis.analysis_result

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Analysis for {res['metadata']['fichier']}"
    slide.placeholders[1].text = f"Date: {res['metadata']['date_analyse']}"

    # Statistics slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Statistics"
    content = slide.placeholders[1]
    for k, v in res['statistiques'].items():
        p = content.text_frame.add_paragraph()
        p.text = f"{k.replace('_', ' ').title()}: {v}"
        p.level = 1
        
    # Alerts slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Alerts & DB Verification"
    content = slide.placeholders[1]
    for k, v in {**res['verification_bd'], **res['alertes']}.items():
        p = content.text_frame.add_paragraph()
        p.text = f"{k.replace('_', ' ').title()}: {v}"
        p.level = 1

    prs.save(buffer)
    buffer.seek(0)
    return buffer